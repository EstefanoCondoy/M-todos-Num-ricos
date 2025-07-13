from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Crear la presentación
prs = Presentation()
slide_width = prs.slide_width

# Función para agregar una diapositiva con título y contenido
def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = "\n".join(content_lines)

# Diapositiva 1: Título
slide_layout = prs.slide_layouts[0]  # Título
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Análisis de Patrones de Consumo de Agua"
slide.placeholders[1].text = "Área: Medio Ambiente\nPresentación de Propuesta Tentativa"

# Diapositiva 2: Descripción del Proyecto
add_slide("Descripción del Proyecto", [
    "• Modelar el consumo de agua en una ciudad durante varios años.",
    "• Identificar patrones de uso, estacionalidades y tendencias.",
    "• Predecir la demanda futura y evaluar escenarios de escasez.",
])

# Diapositiva 3: Objetivo o Meta del Proyecto
add_slide("Objetivo o Meta del Proyecto", [
    "• Crear un modelo matemático del consumo de agua basado en datos reales.",
    "• Estimar el consumo futuro bajo distintas condiciones.",
    "• Simular escenarios críticos ante cambios climáticos.",
])

# Diapositiva 4: Descripción del Dataset
add_slide("Descripción del Dataset", [
    "• Dataset: Consumo de agua en una ciudad de tamaño medio.",
    "• Variables: Year, Month, Temperature, Precipitation, Users, Consumption.",
    "• Fuente: Kaggle - Water Consumption in a Median Size City."
])

# Diapositiva 5: Variables Seleccionadas
add_slide("Variables Seleccionadas", [
    "• Year y Month → Para crear una variable temporal.",
    "• Temperature y Precipitation → Factores climáticos.",
    "• Users → Población que consume agua.",
    "• Consumption → Variable objetivo a modelar.",
])

# Diapositiva 6: Etapa 1 - Limpieza y Exploración de Datos
add_slide("Etapa 1: Limpieza y Exploración", [
    "• Tratamiento de datos nulos o atípicos.",
    "• Fusión de columnas Year y Month.",
    "• Visualización de tendencias y estacionalidades.",
    "• Herramientas: pandas, matplotlib, seaborn.",
])

# Diapositiva 7: Etapa 2 - Análisis de Correlación
add_slide("Etapa 2: Análisis de Correlación", [
    "• Evaluar relaciones entre variables y consumo.",
    "• Gráficos de dispersión y matrices de correlación.",
    "• Métodos numéricos: Cálculo de coeficientes.",
])

# Diapositiva 8: Etapa 3 - Modelado del Consumo (Ajuste de Curvas)
add_slide("Etapa 3: Modelado del Consumo", [
    "• Ajuste de un modelo polinomial por mínimos cuadrados.",
    "• Relación entre tiempo y consumo.",
    "• Base para proyecciones y análisis futuros.",
])

# Diapositiva 9: Etapa 4 - Predicción con Interpolación
add_slide("Etapa 4: Predicción con Interpolación", [
    "• Estimar valores intermedios o futuros.",
    "• Métodos numéricos:",
    "   - Interpolación de Lagrange.",
    "   - Interpolación de Newton.",
])
# Diapositiva 10: Etapa 5 - Simulación de Escenarios
add_slide("Etapa 5: Simulación de Escenarios", [
    "• Simular temperaturas elevadas y baja precipitación.",
    "• Estimar impacto en el consumo.",
    "• Método: uso del modelo ajustado + interpolaciones.",
])

# Diapositiva 11: Resultado Esperado
add_slide("Resultado Esperado", [
    "• Modelo funcional que explique y prediga consumo de agua.",
    "• Proyecciones confiables en escenarios críticos.",
    "• Visualizaciones claras para toma de decisiones.",
])

# Guardar presentación
pptx_path = "C:/Users/Lenovo/Documents/Presentación_Análisis_Patrones_Agua.pptx"

prs.save(pptx_path)
pptx_path


# @autor: Estefano Condoy
# @fecha: 2025-07-01
# @versión: 1.3
# @asignatura: Métodos Numéricos

import numpy as np    # Lo utilizamos solo para crear arreglos y matrices
import time           # Para medir el tiempo de ejecución real
import os             # Para medir el tiempo de CPU del proceso

def Ec_descomposicion_LU_y_resolver(Ec_A, Ec_b):
    """
    Mi código realiza:
    - La descomposición LU de una matriz cuadrada A, es decir A = L · U, tambien resuelve el sistema A · x = b usando LUx = b
    - Y verifica que L · U = A y que A · x ≈ b

    Mis parámetros son:
    - Ec_A: matriz cuadrada (n x n)
    - Ec_b: vector columna (n x 1) o arreglo (n,)

    Lo que va a retonrar es:
    - Matriz L
    - Matriz U
    - Vector solución x
    """

    # Antes de inciar a codificar ponemos el cronómetro y tiempo de CPU para medir el tiempo de ejecución
    Ec_tiempo_inicio = time.time()
    Ec_cpu_inicio = os.times()[0]

    # En el primer paso verificamos que A sea cuadrada y compatible con b
    if Ec_A.shape[0] != Ec_A.shape[1]:
        print("SALIDA: La matriz A no es cuadrada, no se puede descomponer con este método.")
        return None, None, None

    if Ec_A.shape[0] != Ec_b.shape[0]:
        print("SALIDA: Las dimensiones de A y b no son compatibles.")
        return None, None, None

    # En el paso dos hacemos unas copias de trabajo de A y b, es necesario para no modificar las originales
    Ec_A = Ec_A.copy()
    Ec_b = Ec_b.copy()
    Ec_b = Ec_b.flatten()  # <-- Aquí convertimos a vector 1D para evitar problemas
    Ec_n = Ec_A.shape[0]

    # En el paso tres inicializamos L⁻¹ como identidad, U como A
    Ec_L_inv = np.identity(Ec_n)
    Ec_U = Ec_A.copy()

    # En el paso 4 utilizamos la eliliminación Gaussiana
    for Ec_i in range(Ec_n - 1):
        if Ec_U[Ec_i][Ec_i] == 0:
            print("SALIDA: Pivote cero en la posición", Ec_i, ", no se puede continuar.")
            return None, None, None

        for Ec_j in range(Ec_i + 1, Ec_n):
            Ec_mji = Ec_U[Ec_j][Ec_i] / Ec_U[Ec_i][Ec_i]

            Ec_E = np.identity(Ec_n)
            Ec_E[Ec_j][Ec_i] = -Ec_mji

            Ec_U = Ec_E @ Ec_U
            Ec_L_inv = Ec_E @ Ec_L_inv

    # En el paso 5 invertimos L⁻¹ para obtener L
    Ec_L = np.linalg.inv(Ec_L_inv)

    # En el paso 6 ya hacemos la sustitución hacia adelante (Ly = b)
    Ec_y = np.zeros(Ec_n)
    for Ec_i in range(Ec_n):
        Ec_suma = 0
        for Ec_j in range(Ec_i):
            Ec_suma += Ec_L[Ec_i][Ec_j] * Ec_y[Ec_j]
        Ec_y[Ec_i] = (Ec_b[Ec_i] - Ec_suma) / Ec_L[Ec_i][Ec_i]

    # Aqui en el paso 7 tambien hacemos la sustitución hacia atrás (Ux = y)
    Ec_x = np.zeros(Ec_n)
    for Ec_i in range(Ec_n - 1, -1, -1):
        Ec_suma = 0
        for Ec_j in range(Ec_i + 1, Ec_n):
            Ec_suma += Ec_U[Ec_i][Ec_j] * Ec_x[Ec_j]
        Ec_x[Ec_i] = (Ec_y[Ec_i] - Ec_suma) / Ec_U[Ec_i][Ec_i]

    # Verificamos 
    Ec_LU = Ec_L @ Ec_U #El @ es el operador de producto matricial en numpy
    Ec_Ax = Ec_A @ Ec_x

    # Finalizamos el cronometro y tiempo de procesamiento
    Ec_tiempo_fin = time.time()
    Ec_cpu_fin = os.times()[0]

    # Mostramos los resultados y métricas
    print("SALIDA: Ec_Matriz L =\n", Ec_L)
    print("SALIDA: Ec_Matriz U =\n", Ec_U)
    print("PRODUCTO: L · U =\n", Ec_LU)
    print("VECTOR solución x =\n", Ec_x)
    print("VERIFICACIÓN: A · x =\n", Ec_Ax)
    print("VECTOR original b =\n", Ec_b)
    print(f"Tiempo de ejecución real: {Ec_tiempo_fin - Ec_tiempo_inicio:.6f} segundos")
    print(f"Tiempo de procesamiento : {Ec_cpu_fin - Ec_cpu_inicio:.6f} segundos")

    return Ec_L, Ec_U, Ec_x

