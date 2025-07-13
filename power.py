from pptx import Presentation
from pptx.util import Inches, Pt

# Crear presentación
prs = Presentation()

# Función para agregar diapositiva con título y contenido
def add_slide(title, content_lines):
    slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title

    tf = content_placeholder.text_frame
    tf.clear()

    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(5)

# Diapositiva 1: Portada
slide_layout = prs.slide_layouts[0]  # Portada
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Análisis de Patrones de Consumo de Agua"
slide.placeholders[1].text = (
    "Asignatura: Métodos Numéricos\n"
    "Docente: Ayala Tipán Carlos Antonio\n"
    "Paralelo: GR2CC\n"
    "Integrantes: Estefano Condoy y compañeros\n"
    "Fecha: (Colocar fecha de exposición)"
)

# Diapositiva 2: Descripción del Proyecto
add_slide("Descripción del Proyecto", [
    "Modelar el consumo de agua en una ciudad utilizando datos reales.",
    "Identificar patrones de consumo, predecir la demanda futura y evaluar escenarios de escasez.",
    "Enfocado en el área de medio ambiente."
])

# Diapositiva 3: Objetivo del Proyecto
add_slide("Objetivo del Proyecto", [
    "Analizar y predecir el consumo de agua en una ciudad mediana.",
    "Detectar patrones relevantes mediante métodos numéricos.",
    "Simular escenarios futuros considerando variables ambientales y demográficas."
])

# Diapositiva 4: Descripción del Dataset
add_slide("Descripción del Dataset", [
    "Dataset: 'Water Consumption in a Median Size City' (Kaggle).",
    "Registros mensuales de consumo de agua de 2009 a 2016.",
    "Propietario: Empresa AguaH.",
    "Parte de una investigación local."
])

# Diapositiva 5: Variables del Dataset
add_slide("Variables del Dataset", [
    "Fecha (Año/Mes): Para organizar cronológicamente los datos.",
    "Consumo de agua: Variable dependiente principal.",
    "Temperatura promedio y precipitación: Factores climáticos.",
    "Número de habitantes: Relacionado con el crecimiento de la demanda.",
    "Justificación: Variables clave para explicar y predecir el consumo."
])

# Diapositiva 6: Forma Tentativa de Resolución
add_slide("Forma Tentativa de Resolución", [
    "1. Preprocesamiento de datos: limpieza y análisis de estacionalidad.",
    "2. Ajuste de curvas (mínimos cuadrados): modelar tendencias generales.",
    "3. Interpolación (Lagrange, lineal/cuadrática): estimar datos faltantes.",
    "4. Predicción: series de Taylor y splines cúbicos.",
    "5. Simulación de escenarios: método de Newton para ecuaciones planteadas."
])

# Diapositiva 7: Datos adicionales o “Periquitos”
add_slide("Datos adicionales o “Periquitos”", [
    "Según indicaciones del docente, solo se utilizará el dataset oficial.",
    "No se consideran fuentes externas ni variables adicionales.",
    "Se busca garantizar condiciones equitativas para todos los grupos."
])

# Diapositiva 8: Conclusión (Avance Actual)
add_slide("Conclusión (Avance Actual)", [
    "Dataset analizado y comprendido.",
    "Variables seleccionadas.",
    "Métodos numéricos definidos para aplicar.",
    "Proyecto listo para la fase de implementación."
])

# Diapositiva 9: Métodos numéricos que se aplicarán
add_slide("Métodos numéricos que se aplicarán", [
    "Mínimos cuadrados: ajuste de curvas.",
    "Interpolación de Lagrange y lineal: relleno de datos faltantes.",
    "Serie de Taylor y splines cúbicos: predicción.",
    "Método de Newton: resolución de ecuaciones en simulaciones."
])

# Diapositiva 10: Preguntas y Comentarios
add_slide("Preguntas y Comentarios", [
    "¿Dudas o comentarios?"
])

# Guardar presentación
prs.save("Analisis_Patrones_Consumo_Agua.pptx")

